import logging
import subprocess
import csv
import json
from pathlib import Path

import fitz  # pymupdf
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation
import xlrd

logger = logging.getLogger(__name__)


class FinancialReportParser:
    """
    文档解析服务：支持 PDF、Office 和常见文本格式，统一输出 Markdown。
    """

    MARKER_TIMEOUT = 120  # marker 最多等 2 分钟
    SUPPORTED_EXTENSIONS = {
        ".pdf",
        ".txt",
        ".md",
        ".markdown",
        ".csv",
        ".json",
        ".docx",
        ".xlsx",
        ".xls",
        ".pptx",
    }

    # 使用相对于 parser.py 的绝对路径，避免工作目录变动导致数据丢失
    _BASE_DIR = Path(__file__).resolve().parent.parent.parent

    def __init__(self, raw_dir: str | None = None, parsed_dir: str | None = None):
        self.raw_dir = Path(raw_dir) if raw_dir else (self._BASE_DIR / "data" / "raw")
        self.parsed_dir = Path(parsed_dir) if parsed_dir else (self._BASE_DIR / "data" / "parsed")
        self.raw_dir.mkdir(parents=True, exist_ok=True)
        self.parsed_dir.mkdir(parents=True, exist_ok=True)

    # ── 主入口 ──────────────────────────────────────────────
    def is_supported(self, filename: str) -> bool:
        return Path(filename).suffix.lower() in self.SUPPORTED_EXTENSIONS

    def supported_extensions_text(self) -> str:
        return ", ".join(sorted(self.SUPPORTED_EXTENSIONS))

    def parse_file(self, file_path: str) -> str:
        """
        将支持的文件转换为 Markdown，并返回 .md 文件路径。
        """
        source_file = Path(file_path)
        suffix = source_file.suffix.lower()

        if suffix == ".pdf":
            return self.parse_pdf(file_path)
        if suffix in {".txt", ".md", ".markdown"}:
            return self._parse_plain_text(source_file)
        if suffix == ".csv":
            return self._parse_csv(source_file)
        if suffix == ".json":
            return self._parse_json(source_file)
        if suffix == ".docx":
            return self._parse_docx(source_file)
        if suffix == ".xlsx":
            return self._parse_xlsx(source_file)
        if suffix == ".xls":
            return self._parse_xls(source_file)
        if suffix == ".pptx":
            return self._parse_pptx(source_file)

        raise ValueError(f"不支持的文件格式: {suffix}")

    def parse_pdf(self, pdf_path: str, use_marker: bool = False) -> str:
        """
        将 PDF 转换为 Markdown 并返回 .md 文件路径。

        Args:
            pdf_path:   PDF 文件路径
            use_marker: True 则强制使用 marker（高精度，慢）；
                        False 使用 pymupdf 快速提取（默认）
        """
        pdf_file = Path(pdf_path)
        if not pdf_file.exists():
            raise FileNotFoundError(f"未找到 PDF 文件: {pdf_path}")

        output_dir = self.parsed_dir / pdf_file.stem
        output_dir.mkdir(exist_ok=True)
        md_path = output_dir / f"{pdf_file.stem}.md"

        if use_marker:
            return self._parse_with_marker(pdf_file, md_path)

        return self._parse_with_pymupdf(pdf_file, md_path)

    def _make_md_path(self, source_file: Path) -> Path:
        output_dir = self.parsed_dir / source_file.stem
        output_dir.mkdir(exist_ok=True)
        return output_dir / f"{source_file.stem}.md"

    def _write_markdown(self, source_file: Path, content: str) -> str:
        md_path = self._make_md_path(source_file)
        title = f"# {source_file.name}\n\n"
        md_path.write_text(title + content.strip() + "\n", encoding="utf-8")
        return str(md_path)

    def _read_text_with_fallback(self, source_file: Path) -> str:
        for encoding in ("utf-8-sig", "utf-8", "gbk", "gb18030"):
            try:
                return source_file.read_text(encoding=encoding)
            except UnicodeDecodeError:
                continue
        return source_file.read_text(encoding="utf-8", errors="ignore")

    def _parse_plain_text(self, source_file: Path) -> str:
        return self._write_markdown(source_file, self._read_text_with_fallback(source_file))

    def _parse_csv(self, source_file: Path) -> str:
        text = self._read_text_with_fallback(source_file)
        sample = text[:2048]
        try:
            dialect = csv.Sniffer().sniff(sample)
        except csv.Error:
            dialect = csv.excel

        rows = list(csv.reader(text.splitlines(), dialect))
        if not rows:
            return self._write_markdown(source_file, "")

        return self._write_markdown(source_file, self._table_to_markdown(rows))

    def _parse_json(self, source_file: Path) -> str:
        data = json.loads(self._read_text_with_fallback(source_file))
        pretty = json.dumps(data, ensure_ascii=False, indent=2)
        return self._write_markdown(source_file, f"```json\n{pretty}\n```")

    def _parse_docx(self, source_file: Path) -> str:
        doc = DocxDocument(str(source_file))
        parts: list[str] = []

        for paragraph in doc.paragraphs:
            text = paragraph.text.strip()
            if text:
                parts.append(text)

        for table in doc.tables:
            rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
            table_md = self._table_to_markdown(rows)
            if table_md:
                parts.append(table_md)

        return self._write_markdown(source_file, "\n\n".join(parts))

    def _parse_xlsx(self, source_file: Path) -> str:
        workbook = load_workbook(str(source_file), data_only=True, read_only=True)
        parts: list[str] = []

        for sheet in workbook.worksheets:
            rows = []
            for row in sheet.iter_rows(values_only=True):
                values = ["" if value is None else str(value) for value in row]
                if any(value.strip() for value in values):
                    rows.append(values)

            if rows:
                parts.append(f"## Sheet: {sheet.title}\n\n{self._table_to_markdown(rows)}")

        workbook.close()
        return self._write_markdown(source_file, "\n\n".join(parts))

    def _parse_xls(self, source_file: Path) -> str:
        workbook = xlrd.open_workbook(str(source_file))
        parts: list[str] = []

        for sheet in workbook.sheets():
            rows = []
            for row_index in range(sheet.nrows):
                values = [
                    "" if value == "" else str(value)
                    for value in sheet.row_values(row_index)
                ]
                if any(value.strip() for value in values):
                    rows.append(values)

            if rows:
                parts.append(f"## Sheet: {sheet.name}\n\n{self._table_to_markdown(rows)}")

        return self._write_markdown(source_file, "\n\n".join(parts))

    def _parse_pptx(self, source_file: Path) -> str:
        presentation = Presentation(str(source_file))
        parts: list[str] = []

        for index, slide in enumerate(presentation.slides, start=1):
            slide_parts = [f"## Slide {index}"]
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:
                        slide_parts.append(text)
            parts.append("\n\n".join(slide_parts))

        return self._write_markdown(source_file, "\n\n---\n\n".join(parts))

    # ── pymupdf 快速解析 ──────────────────────────────────

    def _parse_with_pymupdf(self, pdf_file: Path, md_path: Path) -> str:
        """用 pymupdf 提取文本，将表格尽量还原为 Markdown 表格。"""
        logger.info("使用快速解析: %s", pdf_file.name)

        doc = fitz.open(str(pdf_file))
        md_parts: list[str] = []

        for page_num, page in enumerate(doc, 1):
            md_parts.append(f"\n\n---\n## Page {page_num}\n")

            # 尝试提取表格（pymupdf 4.x+）
            tables, table_rects = self._extract_tables(page)
            if tables:
                # 有表格时：提取非表格区域文本，避免内容重复
                # 获取所有文本块及其位置信息
                text_blocks = page.get_text("blocks")
                non_table_text = []
                for block in text_blocks:
                    # block: (x0, y0, x1, y1, text, block_no, block_type)
                    if len(block) < 5:
                        continue
                    block_text = block[4].strip() if block[4] else ""
                    if not block_text:
                        continue
                    # 检查文本块是否落在某个表格区域内
                    block_rect = fitz.Rect(block[0], block[1], block[2], block[3])
                    inside_table = any(
                        block_rect.intersects(table_rect)
                        for table_rect in table_rects
                    )
                    if not inside_table:
                        non_table_text.append(block_text)

                if non_table_text:
                    md_parts.append("\n".join(non_table_text))

                for table_md in tables:
                    md_parts.append("\n" + table_md + "\n")
            else:
                text = page.get_text("text").strip()
                if text:
                    md_parts.append(text)

        doc.close()

        md_content = "\n".join(md_parts)
        md_path.write_text(md_content, encoding="utf-8")

        logger.info("快速解析完成: %s", md_path)
        return str(md_path)

    def _extract_tables(self, page) -> tuple[list[str], list[fitz.Rect]]:
        """尝试用 pymupdf 的 find_tables 提取表格并转为 Markdown。

        Returns:
            (markdown_tables, table_rects): Markdown 表格列表和对应的矩形区域列表
        """
        try:
            tabs = page.find_tables()
            if not tabs or len(tabs.tables) == 0:
                return [], []

            results = []
            rects = []
            for table in tabs.tables:
                data = table.extract()
                if not data or len(data) < 2:
                    continue
                results.append(self._table_to_markdown(data))
                rects.append(fitz.Rect(table.bbox))
            return results, rects
        except Exception:
            # pymupdf 旧版本可能没有 find_tables
            return [], []

    @staticmethod
    def _table_to_markdown(rows: list[list[str]]) -> str:
        """将二维数组转为 Markdown 表格。"""
        if not rows:
            return ""

        # 清理 None
        clean = [[cell or "" for cell in row] for row in rows]
        header = clean[0]
        sep = ["-" * max(len(h), 3) for h in header]
        lines = [
            "| " + " | ".join(header) + " |",
            "| " + " | ".join(sep) + " |",
        ]
        for row in clean[1:]:
            # 补齐列数
            while len(row) < len(header):
                row.append("")
            lines.append("| " + " | ".join(row[: len(header)]) + " |")
        return "\n".join(lines)

    # ── marker 高精度解析（后备） ──────────────────────────

    def _parse_with_marker(self, pdf_file: Path, md_path: Path) -> str:
        """调用 marker_single CLI，带超时保护。"""
        logger.info("使用 Marker 高精度解析: %s", pdf_file.name)
        logger.info("超时限制: %ds", self.MARKER_TIMEOUT)

        command = [
            "marker_single",
            str(pdf_file),
            "--output_dir",
            str(self.parsed_dir),
        ]

        try:
            result = subprocess.run(
                command,
                check=True,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                text=True,
                encoding="utf-8",
                timeout=self.MARKER_TIMEOUT,
            )
            logger.info("Marker 解析成功: %s", md_path)
            # 校验 marker 是否确实生成了文件
            if not md_path.exists():
                logger.warning("Marker 返回成功但未生成文件，回退到快速解析")
                return self._parse_with_pymupdf(pdf_file, md_path)
            return str(md_path)

        except subprocess.TimeoutExpired:
            logger.warning("Marker 超时 (%ds)，回退到快速解析", self.MARKER_TIMEOUT)
            return self._parse_with_pymupdf(pdf_file, md_path)

        except (subprocess.CalledProcessError, FileNotFoundError) as e:
            logger.warning("Marker 不可用 (%s)，回退到快速解析", e)
            return self._parse_with_pymupdf(pdf_file, md_path)
